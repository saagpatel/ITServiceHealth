"""Build the Pulse / IT Service Health HDI pitch deck.

10 slides, 16:9, dark theme, single alarm-red accent. No emojis.
Engineering-professional voice. Audience: IT ops leads at HDI SF Bay Area
who run or procure vendor-health tooling.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- design tokens --------------------------------------------------------
BG_PAGE = RGBColor(0x0B, 0x11, 0x20)
SURFACE = RGBColor(0x15, 0x1D, 0x2E)
SURFACE_2 = RGBColor(0x1B, 0x24, 0x36)
BORDER = RGBColor(0x2A, 0x34, 0x46)
TEXT_DISPLAY = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_PRIMARY = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT_ALARM = RGBColor(0xEF, 0x44, 0x44)

HEADER_FONT = "Helvetica Neue"
BODY_FONT = "Helvetica Neue"
MONO_FONT = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- helpers --------------------------------------------------------------
def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        SLIDE_W,
        SLIDE_H,
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # push to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=TEXT_PRIMARY,
    bold=False,
    italic=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    # if text contains newlines, split across runs/paragraphs
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            para = p
        else:
            para = tf.add_paragraph()
            para.alignment = align
        r = para.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill, line_color=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_width is not None:
            s.line.width = line_width
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, *, fill, line_color=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_width is not None:
            s.line.width = line_width
    s.shadow.inherit = False
    # set small corner radius — adjust adjustment value
    try:
        s.adjustments[0] = 0.08
    except (IndexError, ValueError):
        pass
    return s


def add_hairline(slide, x, y, w, color=BORDER):
    return add_rect(slide, x, y, w, Emu(12700), fill=color)  # ~1px


def slide_header(slide, eyebrow: str, title: str):
    """Consistent slide header — eyebrow + title + alarm-red rule."""
    add_text(
        slide,
        eyebrow,
        Inches(0.6),
        Inches(0.45),
        Inches(10),
        Inches(0.3),
        size=11,
        color=ACCENT_ALARM,
        bold=True,
        font=MONO_FONT,
    )
    add_text(
        slide,
        title,
        Inches(0.6),
        Inches(0.75),
        Inches(12),
        Inches(0.9),
        size=36,
        color=TEXT_DISPLAY,
        bold=True,
        font=HEADER_FONT,
    )


def page_number(slide, n: int, total: int = 10):
    add_text(
        slide,
        f"{n:02d} / {total:02d}",
        Inches(12.3),
        Inches(7.05),
        Inches(1),
        Inches(0.3),
        size=9,
        color=TEXT_MUTED,
        font=MONO_FONT,
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        "Pulse · HDI SF Bay Area",
        Inches(0.6),
        Inches(7.05),
        Inches(6),
        Inches(0.3),
        size=9,
        color=TEXT_MUTED,
        font=MONO_FONT,
    )


# --- presentation setup ---------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# --- Slide 1: title -------------------------------------------------------
def slide_title():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)

    # Left accent bar
    add_rect(s, Inches(0.6), Inches(2.9), Inches(0.15), Inches(1.2), fill=ACCENT_ALARM)

    add_text(
        s,
        "HDI · SF Bay Area · Vendor-health tooling",
        Inches(0.9),
        Inches(2.9),
        Inches(11),
        Inches(0.3),
        size=12,
        color=ACCENT_ALARM,
        bold=True,
        font=MONO_FONT,
    )
    add_text(
        s,
        "Pulse",
        Inches(0.9),
        Inches(3.2),
        Inches(12),
        Inches(1.3),
        size=88,
        color=TEXT_DISPLAY,
        bold=True,
        font=HEADER_FONT,
    )
    add_text(
        s,
        "Production-grade vendor-health monitoring for IT ops.",
        Inches(0.9),
        Inches(4.5),
        Inches(12),
        Inches(0.6),
        size=22,
        color=TEXT_SECONDARY,
        font=HEADER_FONT,
    )
    add_text(
        s,
        "Thirty SaaS vendors. Sixty-second cadence. Two hundred seventy-six tests.\n"
        "Built in an enterprise IT environment as a Platform Engineer transition proof point.",
        Inches(0.9),
        Inches(5.2),
        Inches(12),
        Inches(1.0),
        size=14,
        color=TEXT_MUTED,
        font=BODY_FONT,
    )

    add_text(
        s,
        "Saagar Patel · Platform Engineering",
        Inches(0.6),
        Inches(6.6),
        Inches(8),
        Inches(0.3),
        size=11,
        color=TEXT_MUTED,
        font=MONO_FONT,
    )
    add_text(
        s,
        "github.com/saagpatel/ITServiceHealth · MIT",
        Inches(5),
        Inches(6.6),
        Inches(8),
        Inches(0.3),
        size=11,
        color=TEXT_MUTED,
        font=MONO_FONT,
        align=PP_ALIGN.RIGHT,
    )


# --- Slide 2: the problem -------------------------------------------------
def slide_problem():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(
        s, "01 · PROBLEM", "The vendor-health signal is broken at the operator's desk."
    )

    # four callout cards in a 2x2 grid
    cards = [
        (
            "30+ SaaS vendors",
            "Identity, productivity, collaboration, engineering, HR, finance, sales, marketing, support. "
            "Every outage reaches the help desk before the status page does.",
        ),
        (
            "Status pages lag reality",
            "Vendors mark themselves green while users flood the IT help channel. RSS is stale. JSON APIs don't "
            "agree on a schema. Manual updates go stale within the hour.",
        ),
        (
            "Slack becomes the war room",
            "No dedup. No severity routing. Alerts arrive on the same channel as brownouts, flapping "
            "pollers, and vendor maintenance windows no one acknowledged.",
        ),
        (
            "Leadership asks: is it us?",
            "Exec walks up to the big screen and wants one number. The NOC wall shows 80 service tiles "
            "and a timeline. That is not the meter a director reads.",
        ),
    ]
    col_w, row_h = Inches(6.0), Inches(2.4)
    x0, y0 = Inches(0.6), Inches(1.9)
    for i, (head, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = x0 + col * (col_w + Inches(0.25))
        cy = y0 + row * (row_h + Inches(0.25))
        add_rounded(
            s, cx, cy, col_w, row_h, fill=SURFACE, line_color=BORDER, line_width=Pt(0.5)
        )
        # left alarm bar
        add_rect(
            s,
            cx + Inches(0.25),
            cy + Inches(0.3),
            Inches(0.05),
            Inches(0.5),
            fill=ACCENT_ALARM,
        )
        add_text(
            s,
            head,
            cx + Inches(0.5),
            cy + Inches(0.25),
            col_w - Inches(0.7),
            Inches(0.5),
            size=20,
            color=TEXT_DISPLAY,
            bold=True,
            font=HEADER_FONT,
        )
        add_text(
            s,
            body,
            cx + Inches(0.5),
            cy + Inches(0.85),
            col_w - Inches(0.7),
            Inches(1.5),
            size=13,
            color=TEXT_SECONDARY,
            font=BODY_FONT,
        )

    page_number(s, 2)


# --- Slide 3: what we built ----------------------------------------------
def slide_what():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(
        s, "02 · WHAT WE BUILT", "One dashboard. Two views. Production-graded."
    )

    # left column: big statement
    add_text(
        s,
        "Pulse polls every vendor, normalizes five states, detects changes, "
        "writes audit events, and posts one clean Slack alert per real "
        "incident. Two views off the same data: Executive for the "
        "conference room, Engineer for the triage desk.",
        Inches(0.6),
        Inches(1.95),
        Inches(7.3),
        Inches(3.5),
        size=17,
        color=TEXT_PRIMARY,
        font=HEADER_FONT,
    )

    # right column: shipped feature chip list
    chip_x = Inches(8.3)
    chip_y = Inches(1.95)
    chip_w = Inches(4.45)
    chips = [
        (
            "v1",
            "Poll loop · 5-state normalizer · Slack Block Kit · React UI · dep graph · timeline · SLA",
        ),
        (
            "v2 · phase 0-1",
            "Bearer-token admin auth · stamina retries · purgatory circuit breakers · poller_health · unknown-on-blind",
        ),
        (
            "v2 · phase 2",
            "Flap suppression · dedup window · tier routing · dependency correlation · maintenance windows",
        ),
        (
            "v2 · phase 3",
            "structlog JSON · Prometheus /metrics · Sentry · Healthchecks.io dead-man's switch",
        ),
        (
            "v2 · phase 4",
            "aiosqlite pool · Litestream streaming · daily VACUUM INTO · retention · WAL checkpointing",
        ),
        (
            "v2 · phase 5-6",
            "TanStack-style polling · Executive/Engineer toggle · PWA · a11y · CI · Caddy · Keychain",
        ),
    ]
    row_h = Inches(0.72)
    for i, (tag, body) in enumerate(chips):
        ry = chip_y + i * row_h
        add_rect(
            s,
            chip_x,
            ry + Inches(0.08),
            Inches(0.08),
            Inches(0.5),
            fill=ACCENT_ALARM if "v2" in tag else TEXT_MUTED,
        )
        add_text(
            s,
            tag,
            chip_x + Inches(0.25),
            ry + Inches(0.05),
            Inches(1.4),
            Inches(0.3),
            size=10,
            color=ACCENT_ALARM if "v2" in tag else TEXT_MUTED,
            bold=True,
            font=MONO_FONT,
        )
        add_text(
            s,
            body,
            chip_x + Inches(0.25),
            ry + Inches(0.32),
            chip_w - Inches(0.3),
            Inches(0.4),
            size=11,
            color=TEXT_SECONDARY,
            font=BODY_FONT,
        )

    # bottom stat row
    stats_y = Inches(6.1)
    stat_defs = [
        ("30", "SaaS vendors"),
        ("60 s", "poll cadence"),
        ("276", "tests passing"),
        ("~1 s", "RPO · Litestream"),
    ]
    stat_w = Inches(3.1)
    for i, (num, label) in enumerate(stat_defs):
        sx = Inches(0.6) + i * (stat_w + Inches(0.05))
        add_text(
            s,
            num,
            sx,
            stats_y,
            stat_w,
            Inches(0.55),
            size=28,
            color=TEXT_DISPLAY,
            bold=True,
            font=HEADER_FONT,
        )
        add_text(
            s,
            label,
            sx,
            stats_y + Inches(0.55),
            stat_w,
            Inches(0.3),
            size=10,
            color=TEXT_MUTED,
            font=MONO_FONT,
        )

    page_number(s, 3)


# --- Slide 4: executive view screenshot ----------------------------------
def slide_exec():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(
        s,
        "03 · EXECUTIVE VIEW",
        "One panel. Three meters. Read it from the back of the room.",
    )

    # Mockup is 1920x1080 (16:9). Fit by height so it doesn't overflow the
    # slide — height 5.0 in → width 5.0 * 16/9 = 8.89 in. Center horizontally.
    img_path = "docs/executive-view-redesign/screenshots/exec-major.png"
    img_w = Inches(8.89)
    img_h = Inches(5.0)
    img_x = (SLIDE_W - img_w) / 2
    s.shapes.add_picture(img_path, img_x, Inches(1.85), width=img_w, height=img_h)

    add_text(
        s,
        "Executive view · rendered at 1920x1080 · 2 active incidents",
        img_x,
        Inches(6.9),
        img_w,
        Inches(0.22),
        size=9,
        color=TEXT_MUTED,
        font=MONO_FONT,
        align=PP_ALIGN.CENTER,
    )

    page_number(s, 4)


# --- Slide 5: architecture diagram ---------------------------------------
def slide_arch():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(
        s,
        "04 · ARCHITECTURE",
        "Vendor pages → resilient poll → alert hygiene → dashboard.",
    )

    # Architecture PNG is 11x8.5 landscape (ratio 1.294). Fit by height
    # so the diagram stays on-slide. height 5.0 → width ~6.47; that's narrow
    # so size up to fill: height 5.0, width = 5.0 * 11/8.5 = 6.47. Widen by
    # stretching horizontally because the source has breathable whitespace.
    img_w = Inches(10.5)
    img_h = img_w * 8.5 / 11  # ~8.11 in — still too tall
    # Constrain by height instead.
    img_h = Inches(5.0)
    img_w = img_h * 11 / 8.5
    img_x = (SLIDE_W - img_w) / 2
    s.shapes.add_picture(
        "docs/architecture-diagram/architecture.png",
        img_x,
        Inches(1.85),
        width=img_w,
        height=img_h,
    )

    add_text(
        s,
        "Landscape architecture · resilience + alert + observability lanes",
        img_x,
        Inches(6.9),
        img_w,
        Inches(0.22),
        size=9,
        color=TEXT_MUTED,
        font=MONO_FONT,
        align=PP_ALIGN.CENTER,
    )

    page_number(s, 5)


# --- Slide 6: resilience gates --------------------------------------------
def slide_resilience():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(s, "05 · RESILIENCE", "Vendor APIs misbehave. The poller does not.")

    # two large cards side by side — stamina and purgatory
    card_y = Inches(1.95)
    card_h = Inches(4.4)
    card_w = Inches(6.0)

    # stamina
    cx = Inches(0.6)
    add_rounded(
        s,
        cx,
        card_y,
        card_w,
        card_h,
        fill=SURFACE,
        line_color=BORDER,
        line_width=Pt(0.5),
    )
    add_text(
        s,
        "stamina",
        cx + Inches(0.4),
        card_y + Inches(0.35),
        card_w - Inches(0.8),
        Inches(0.5),
        size=28,
        color=TEXT_DISPLAY,
        bold=True,
        font=HEADER_FONT,
    )
    add_text(
        s,
        "Exponential backoff with jitter on every outbound call.",
        cx + Inches(0.4),
        card_y + Inches(1.0),
        card_w - Inches(0.8),
        Inches(0.5),
        size=14,
        color=TEXT_SECONDARY,
        font=BODY_FONT,
    )
    stamina_lines = [
        "retry budget · 3 attempts with exponential base 0.5 s",
        "jitter · ±250 ms so we don't synchronise across services",
        "timeout · 10 s read · 5 s connect · no blocking I/O",
        "observability · every retry labeled with vendor + attempt",
    ]
    for i, line in enumerate(stamina_lines):
        add_text(
            s,
            "· " + line,
            cx + Inches(0.4),
            card_y + Inches(1.9) + i * Inches(0.55),
            card_w - Inches(0.8),
            Inches(0.4),
            size=12,
            color=TEXT_PRIMARY,
            font=BODY_FONT,
        )

    # purgatory
    cx = Inches(6.75)
    add_rounded(
        s,
        cx,
        card_y,
        card_w,
        card_h,
        fill=SURFACE,
        line_color=ACCENT_ALARM,
        line_width=Pt(1.25),
    )
    add_text(
        s,
        "purgatory",
        cx + Inches(0.4),
        card_y + Inches(0.35),
        card_w - Inches(0.8),
        Inches(0.5),
        size=28,
        color=TEXT_DISPLAY,
        bold=True,
        font=HEADER_FONT,
    )
    add_text(
        s,
        "Per-host circuit breaker. Blind is not operational.",
        cx + Inches(0.4),
        card_y + Inches(1.0),
        card_w - Inches(0.8),
        Inches(0.5),
        size=14,
        color=TEXT_SECONDARY,
        font=BODY_FONT,
    )
    purg_lines = [
        "3 consecutive failures · breaker opens for 300 s",
        "half-open probe · one attempt after TTL before closing",
        "poller_health flips · service renders as unknown, not operational",
        "Slack notification · separate channel from incident alerts",
    ]
    for i, line in enumerate(purg_lines):
        add_text(
            s,
            "· " + line,
            cx + Inches(0.4),
            card_y + Inches(1.9) + i * Inches(0.55),
            card_w - Inches(0.8),
            Inches(0.4),
            size=12,
            color=TEXT_PRIMARY,
            font=BODY_FONT,
        )

    add_text(
        s,
        "Rule: if we cannot see the vendor, we do not claim they are up. "
        "This is the single dashboard bug you cannot ship with.",
        Inches(0.6),
        Inches(6.55),
        Inches(12.1),
        Inches(0.4),
        size=13,
        color=ACCENT_ALARM,
        italic=True,
        font=HEADER_FONT,
    )

    page_number(s, 6)


# --- Slide 7: alert hygiene -----------------------------------------------
def slide_alerts():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(
        s, "06 · ALERT HYGIENE", "One message per real incident. No one wakes up twice."
    )

    # Five numbered rows
    rows = [
        (
            "Flap suppression",
            "3 consecutive confirming polls before worsening alerts fire. "
            "Minimum state duration of 600 s. Recovery requires 2 consecutive operational polls.",
        ),
        (
            "Deduplication",
            "Alert dedup keyed on vendor_incident_id, not message text. "
            "One-day window by default. Same incident cannot re-page on every poll.",
        ),
        (
            "Tier routing",
            "Critical vs informational. Tier picks destination channel and Slack priority. "
            "Degraded-brownout and major-outage do not share a notification surface.",
        ),
        (
            "Dependency correlation",
            "When the identity provider goes down, we send one aggregated upstream alert instead of "
            "cascading alerts for every dependent service. Threshold is configurable.",
        ),
        (
            "Maintenance windows",
            "First-class DB table. Scheduled vendor windows suppress alerts for the duration. "
            "Auto-populated from vendor feeds; manual windows supported.",
        ),
    ]
    row_h = Inches(0.88)
    y0 = Inches(1.95)
    for i, (title, body) in enumerate(rows):
        ry = y0 + i * row_h
        # number chip
        add_rounded(s, Inches(0.6), ry, Inches(0.7), Inches(0.7), fill=ACCENT_ALARM)
        add_text(
            s,
            f"{i + 1:02d}",
            Inches(0.6),
            ry + Inches(0.1),
            Inches(0.7),
            Inches(0.5),
            size=18,
            color=TEXT_DISPLAY,
            bold=True,
            font=MONO_FONT,
            align=PP_ALIGN.CENTER,
        )
        # title + body
        add_text(
            s,
            title,
            Inches(1.55),
            ry,
            Inches(4.2),
            Inches(0.4),
            size=16,
            color=TEXT_DISPLAY,
            bold=True,
            font=HEADER_FONT,
        )
        add_text(
            s,
            body,
            Inches(5.75),
            ry,
            Inches(7.1),
            Inches(0.75),
            size=12,
            color=TEXT_SECONDARY,
            font=BODY_FONT,
        )

    add_text(
        s,
        "Env-tunable: ALERT_CONFIRM_THRESHOLD_POLLS · "
        "ALERT_DEDUP_WINDOW_SECONDS · DEPENDENCY_CORRELATION_THRESHOLD",
        Inches(0.6),
        Inches(6.55),
        Inches(12.1),
        Inches(0.4),
        size=10,
        color=TEXT_MUTED,
        font=MONO_FONT,
    )

    page_number(s, 7)


# --- Slide 8: observability -----------------------------------------------
def slide_obs():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(
        s,
        "07 · OBSERVABILITY",
        "If the dashboard were down, I would know in 30 seconds.",
    )

    # four observability cards in one row
    cards = [
        (
            "structlog",
            "JSON logs",
            "Every request, poll, and alert carries a correlation id. "
            "WatchedFileHandler survives logrotate.",
        ),
        (
            "Prometheus",
            "/metrics exposition",
            "Poll latency, breaker state, alert dispatches, SLA percentage per service. "
            "Scrapeable by any text-format collector.",
        ),
        (
            "Sentry",
            "Exception capture",
            "Unhandled exceptions reported with release tag. "
            "Traces optional. Default sample rate zero.",
        ),
        (
            "Healthchecks.io",
            "Dead-man's switch",
            "30-second heartbeat. /healthz returns 503 past 120 s. "
            "External observer notices silence even if monitoring itself is dead.",
        ),
    ]
    card_w = Inches(3.0)
    card_h = Inches(4.4)
    gap = Inches(0.12)
    y0 = Inches(1.95)
    x0 = Inches(0.6)
    for i, (name, role, body) in enumerate(cards):
        cx = x0 + i * (card_w + gap)
        add_rounded(
            s,
            cx,
            y0,
            card_w,
            card_h,
            fill=SURFACE,
            line_color=BORDER,
            line_width=Pt(0.5),
        )
        add_text(
            s,
            name,
            cx + Inches(0.3),
            y0 + Inches(0.4),
            card_w - Inches(0.5),
            Inches(0.6),
            size=19,
            color=TEXT_DISPLAY,
            bold=True,
            font=HEADER_FONT,
        )
        add_text(
            s,
            role,
            cx + Inches(0.4),
            y0 + Inches(1.05),
            card_w - Inches(0.8),
            Inches(0.4),
            size=11,
            color=ACCENT_ALARM,
            bold=True,
            font=MONO_FONT,
        )
        add_text(
            s,
            body,
            cx + Inches(0.4),
            y0 + Inches(1.55),
            card_w - Inches(0.8),
            Inches(2.7),
            size=12,
            color=TEXT_SECONDARY,
            font=BODY_FONT,
        )

    add_text(
        s,
        "Rule: every layer reports independently. The poller, the API, "
        "the writer, and the browser all answer the question is this alive?",
        Inches(0.6),
        Inches(6.55),
        Inches(12.1),
        Inches(0.4),
        size=13,
        color=TEXT_PRIMARY,
        italic=True,
        font=HEADER_FONT,
    )

    page_number(s, 8)


# --- Slide 9: results -----------------------------------------------------
def slide_results():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)
    slide_header(
        s, "08 · RESULTS", "Shipped. In tree. Production-graded. Boring by design."
    )

    # left column: big stats
    stats = [
        ("7", "production phases shipped"),
        ("276", "tests passing"),
        ("~1 s", "RPO · Litestream streaming"),
        ("< 30 s", "p50 vendor-status detection"),
    ]
    y0 = Inches(1.95)
    for i, (num, label) in enumerate(stats):
        sx = Inches(0.6)
        sy = y0 + i * Inches(1.1)
        add_text(
            s,
            num,
            sx,
            sy,
            Inches(3.0),
            Inches(0.75),
            size=52,
            color=TEXT_DISPLAY,
            bold=True,
            font=HEADER_FONT,
        )
        add_text(
            s,
            label,
            sx + Inches(3.1),
            sy + Inches(0.25),
            Inches(4.0),
            Inches(0.4),
            size=13,
            color=TEXT_MUTED,
            font=MONO_FONT,
        )

    # right column: before/after
    rcx = Inches(7.8)
    rcw = Inches(5.0)
    add_rounded(
        s,
        rcx,
        y0,
        rcw,
        Inches(4.4),
        fill=SURFACE,
        line_color=BORDER,
        line_width=Pt(0.5),
    )
    add_text(
        s,
        "Before → After",
        rcx + Inches(0.4),
        y0 + Inches(0.3),
        rcw - Inches(0.8),
        Inches(0.5),
        size=18,
        color=TEXT_DISPLAY,
        bold=True,
        font=HEADER_FONT,
    )
    ba = [
        ("Detection", "user reports in Slack", "60 s poll + dedup alert"),
        ("Reliability", "vendor page = truth", "blind = unknown, not green"),
        ("Exec read", "80-tile NOC grid", "one panel · three meters"),
        ("Alerts", "one per poll", "one per real incident"),
        ("Backups", "none", "~1 s RPO + daily snapshot"),
    ]
    for i, (k, v1, v2) in enumerate(ba):
        ly = y0 + Inches(1.1) + i * Inches(0.62)
        add_text(
            s,
            k,
            rcx + Inches(0.4),
            ly,
            Inches(1.2),
            Inches(0.3),
            size=11,
            color=TEXT_MUTED,
            font=MONO_FONT,
            bold=True,
        )
        add_text(
            s,
            v1,
            rcx + Inches(1.55),
            ly,
            Inches(1.65),
            Inches(0.3),
            size=11,
            color=TEXT_SECONDARY,
            font=BODY_FONT,
        )
        add_text(
            s,
            "→",
            rcx + Inches(3.1),
            ly,
            Inches(0.3),
            Inches(0.3),
            size=11,
            color=ACCENT_ALARM,
            bold=True,
            font=MONO_FONT,
        )
        add_text(
            s,
            v2,
            rcx + Inches(3.4),
            ly,
            Inches(1.4),
            Inches(0.3),
            size=11,
            color=TEXT_PRIMARY,
            font=BODY_FONT,
            bold=True,
        )

    add_text(
        s,
        "Not shipped: LLM summarization, Splunk / JSM / ThousandEyes / Datadog. Deferred on purpose.",
        Inches(0.6),
        Inches(6.55),
        Inches(12.1),
        Inches(0.4),
        size=11,
        color=TEXT_MUTED,
        italic=True,
        font=BODY_FONT,
    )

    page_number(s, 9)


# --- Slide 10: platform engineer lens ------------------------------------
def slide_close():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, BG_PAGE)

    # large statement — closing mirrors the title slide
    add_rect(s, Inches(0.6), Inches(2.2), Inches(0.15), Inches(1.2), fill=ACCENT_ALARM)

    add_text(
        s,
        "09 · PLATFORM ENGINEER LENS",
        Inches(0.9),
        Inches(2.2),
        Inches(12),
        Inches(0.3),
        size=12,
        color=ACCENT_ALARM,
        bold=True,
        font=MONO_FONT,
    )

    add_text(
        s,
        "This is what the job looks like\nfrom this side of the help-desk.",
        Inches(0.9),
        Inches(2.55),
        Inches(12),
        Inches(1.8),
        size=40,
        color=TEXT_DISPLAY,
        bold=True,
        font=HEADER_FONT,
    )

    body = (
        "Pulse is not a side project. It is the shape of the work: listen to vendor "
        "signal end-to-end, design resilience gates before you write business "
        "logic, refuse to claim green when you are blind, put one meter in the "
        "room that a director can read. IT support told me where the pain was. "
        "Platform engineering is what I did with it."
    )
    add_text(
        s,
        body,
        Inches(0.9),
        Inches(4.4),
        Inches(11.8),
        Inches(1.8),
        size=15,
        color=TEXT_SECONDARY,
        font=BODY_FONT,
    )

    # bottom band — what to take away
    add_hairline(s, Inches(0.6), Inches(6.35), Inches(12.1), color=BORDER)

    takeaways = [
        ("If you operate", "Steal the five alert-hygiene layers."),
        ("If you procure", "Ask vendors for these five guarantees, not feature lists."),
        ("If you build", "Clone it. MIT. github.com/saagpatel/ITServiceHealth."),
    ]
    tw = Inches(4.0)
    for i, (tag, msg) in enumerate(takeaways):
        tx = Inches(0.6) + i * (tw + Inches(0.1))
        add_text(
            s,
            tag,
            tx,
            Inches(6.55),
            tw,
            Inches(0.3),
            size=10,
            color=ACCENT_ALARM,
            bold=True,
            font=MONO_FONT,
        )
        add_text(
            s,
            msg,
            tx,
            Inches(6.8),
            tw,
            Inches(0.45),
            size=13,
            color=TEXT_PRIMARY,
            bold=True,
            font=HEADER_FONT,
        )

    add_text(
        s,
        "10 / 10",
        Inches(12.3),
        Inches(7.05),
        Inches(1),
        Inches(0.3),
        size=9,
        color=TEXT_MUTED,
        font=MONO_FONT,
        align=PP_ALIGN.RIGHT,
    )


# --- render ---------------------------------------------------------------
slide_title()
slide_problem()
slide_what()
slide_exec()
slide_arch()
slide_resilience()
slide_alerts()
slide_obs()
slide_results()
slide_close()

out = "docs/pitch-deck/pulse-hdi.pptx"
prs.save(out)
print(f"wrote {out} · {len(prs.slides)} slides")
